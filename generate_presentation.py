"""
Générateur de Présentation PowerPoint Professionnelle
Interface de Comparaison PACS - Orthanc vs DCM4CHEE
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def add_title_slide(prs, title, subtitle):
    """Slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background gradient (simulation avec rectangle)
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(102, 126, 234)  # Primary color
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(0.5)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Présenté par: Votre Nom | Janvier 2026"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(14)
    footer_p.font.color.rgb = RGBColor(255, 255, 255)
    footer_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Slide avec titre et contenu bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Title
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    # Content
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(12)
    
    return slide

def add_table_slide(prs, title, data):
    """Slide avec tableau"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(102, 126, 234)
    
    # Table
    rows = len(data)
    cols = len(data[0])
    
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    for col_idx, cell_text in enumerate(data[0]):
        cell = table.rows[0].cells[col_idx]
        cell.text = str(cell_text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(102, 126, 234)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx in range(1, rows):
        for col_idx, cell_text in enumerate(data[row_idx]):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if col_idx == 0:
                cell.text_frame.paragraphs[0].font.bold = True
    
    return slide

def add_chart_slide(prs, title, categories, values1, values2, label1, label2):
    """Slide avec graphique comparatif"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(102, 126, 234)
    
    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(label1, values1)
    chart_data.add_series(label2, values2)
    
    # Add chart
    x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    return slide

def add_architecture_slide(prs):
    """Slide architecture technique"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture Technique"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(102, 126, 234)
    
    # Frontend box
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(1.2), Inches(6), Inches(0.8)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(16, 185, 129)
    shape1.text = "Frontend React 18 + Vite\nDashboard + OHIF Viewer"
    shape1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape1.text_frame.paragraphs[0].font.size = Pt(14)
    shape1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Backend box
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2.5), Inches(6), Inches(0.8)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(102, 126, 234)
    shape2.text = "Backend FastAPI (Python 3.13)\nOrchestration + Anonymisation"
    shape2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape2.text_frame.paragraphs[0].font.size = Pt(14)
    shape2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # PACS boxes
    shape3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(3.8), Inches(2.3), Inches(0.7)
    )
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(245, 158, 11)
    shape3.text = "Orthanc :8042"
    shape3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape3.text_frame.paragraphs[0].font.size = Pt(14)
    shape3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    shape4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.2), Inches(3.8), Inches(2.3), Inches(0.7)
    )
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = RGBColor(245, 158, 11)
    shape4.text = "DCM4CHEE :8080"
    shape4.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape4.text_frame.paragraphs[0].font.size = Pt(14)
    shape4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    shape5 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(3.8), Inches(2.3), Inches(0.7)
    )
    shape5.fill.solid()
    shape5.fill.fore_color.rgb = RGBColor(245, 158, 11)
    shape5.text = "XNAT :8090"
    shape5.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape5.text_frame.paragraphs[0].font.size = Pt(14)
    shape5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Database box
    shape6 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(5), Inches(5), Inches(0.7)
    )
    shape6.fill.solid()
    shape6.fill.fore_color.rgb = RGBColor(59, 130, 246)
    shape6.text = "PostgreSQL 15 - Logs + Audit + Anonymisation"
    shape6.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape6.text_frame.paragraphs[0].font.size = Pt(14)
    shape6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add arrows (lines)
    connector1 = slide.shapes.add_connector(1, Inches(5), Inches(2), Inches(5), Inches(2.5))
    connector1.line.color.rgb = RGBColor(0, 0, 0)
    connector1.line.width = Pt(2)
    
    return slide

def add_objectives_slide(prs):
    """Slide objectifs structuré en 3 catégories"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "3️⃣ Objectifs du Projet"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(102, 126, 234)
    
    # Section 1: Objectifs Techniques
    section1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(3), Inches(3.5)
    )
    section1.fill.solid()
    section1.fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    tf1 = section1.text_frame
    tf1.clear()
    tf1.word_wrap = True
    
    p1_title = tf1.add_paragraph()
    p1_title.text = "🔧 Objectifs Techniques"
    p1_title.font.size = Pt(16)
    p1_title.font.bold = True
    p1_title.font.color.rgb = RGBColor(255, 255, 255)
    p1_title.space_after = Pt(12)
    
    objectives_tech = [
        "• Comparer Orthanc vs DCM4CHEE",
        "• Backend FastAPI performant",
        "• Architecture microservices",
        "• APIs REST sécurisées",
        "• Docker containerisé"
    ]
    
    for obj in objectives_tech:
        p = tf1.add_paragraph()
        p.text = obj
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(6)
    
    # Section 2: Objectifs Fonctionnels
    section2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.7), Inches(1.2), Inches(3), Inches(3.5)
    )
    section2.fill.solid()
    section2.fill.fore_color.rgb = RGBColor(16, 185, 129)
    
    tf2 = section2.text_frame
    tf2.clear()
    tf2.word_wrap = True
    
    p2_title = tf2.add_paragraph()
    p2_title.text = "🎯 Objectifs Fonctionnels"
    p2_title.font.size = Pt(16)
    p2_title.font.bold = True
    p2_title.font.color.rgb = RGBColor(255, 255, 255)
    p2_title.space_after = Pt(12)
    
    objectives_func = [
        "• Dashboard comparatif React",
        "• Viewer OHIF intégré",
        "• Extraction RT-STRUCT 3D",
        "• Interface intuitive",
        "• Temps réel (30s)"
    ]
    
    for obj in objectives_func:
        p = tf2.add_paragraph()
        p.text = obj
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(6)
    
    # Section 3: Objectifs Conformité
    section3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.9), Inches(1.2), Inches(3), Inches(3.5)
    )
    section3.fill.solid()
    section3.fill.fore_color.rgb = RGBColor(245, 158, 11)
    
    tf3 = section3.text_frame
    tf3.clear()
    tf3.word_wrap = True
    
    p3_title = tf3.add_paragraph()
    p3_title.text = "🔒 Objectifs Conformité"
    p3_title.font.size = Pt(16)
    p3_title.font.bold = True
    p3_title.font.color.rgb = RGBColor(255, 255, 255)
    p3_title.space_after = Pt(12)
    
    objectives_conf = [
        "• Anonymisation RGPD/HIPAA",
        "• Intégration XNAT recherche",
        "• Traçabilité complète",
        "• Tests 87% couverture",
        "• Documentation exhaustive"
    ]
    
    for obj in objectives_conf:
        p = tf3.add_paragraph()
        p.text = obj
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(6)
    
    # Bottom summary box
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.2), Inches(7), Inches(1.2)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(71, 85, 105)
    
    tf_summary = summary_box.text_frame
    tf_summary.clear()
    tf_summary.word_wrap = True
    
    p_summary_title = tf_summary.add_paragraph()
    p_summary_title.text = "🎯 Objectif Global"
    p_summary_title.font.size = Pt(14)
    p_summary_title.font.bold = True
    p_summary_title.font.color.rgb = RGBColor(255, 255, 255)
    p_summary_title.alignment = PP_ALIGN.CENTER
    
    p_summary = tf_summary.add_paragraph()
    p_summary.text = "Créer une plateforme complète de comparaison PACS avec anonymisation pour la recherche clinique"
    p_summary.font.size = Pt(13)
    p_summary.font.color.rgb = RGBColor(255, 255, 255)
    p_summary.alignment = PP_ALIGN.CENTER
    
    return slide

def add_stats_slide(prs, title, stats):
    """Slide avec statistiques visuelles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(102, 126, 234)
    
    # Stats boxes (3 columns)
    colors = [
        RGBColor(102, 126, 234),
        RGBColor(16, 185, 129),
        RGBColor(245, 158, 11)
    ]
    
    for idx, (label, value) in enumerate(stats.items()):
        col = idx % 3
        row = idx // 3
        
        left = Inches(1 + col * 2.8)
        top = Inches(2 + row * 1.8)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(2.3), Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[col]
        
        # Value
        text_frame = shape.text_frame
        text_frame.clear()
        
        p1 = text_frame.add_paragraph()
        p1.text = str(value)
        p1.font.size = Pt(32)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        
        # Label
        p2 = text_frame.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_presentation():
    """Génère la présentation complète"""
    print("🎨 Génération de la présentation PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Titre
    print("  📄 Slide 1/25: Titre")
    add_title_slide(
        prs,
        "🏥 Interface de Comparaison PACS",
        "Analyse Comparative Orthanc vs DCM4CHEE\nIntégration XNAT et Extraction RT-STRUCT"
    )
    
    # Slide 2: Table des matières
    print("  📄 Slide 2/25: Table des matières")
    add_content_slide(prs, "📋 Table des Matières", [
        "1. Introduction et Contexte",
        "2. Problématique",
        "3. Objectifs du Projet",
        "4. Architecture Technique",
        "5. Technologies et Stack",
        "6. Fonctionnalités Principales",
        "7. Dashboard Comparatif",
        "8. Visualisation DICOM",
        "9. Anonymisation RGPD/HIPAA",
        "10. Extraction RT-STRUCT",
        "11. Résultats et Métriques",
        "12. Tests et Qualité",
        "13. Démonstration",
        "14. Perspectives",
        "15. Conclusion"
    ])
    
    # Slide 3: Introduction
    print("  📄 Slide 3/25: Introduction")
    add_content_slide(prs, "1️⃣ Introduction et Contexte", [
        "🏥 PACS (Picture Archiving and Communication System)",
        "📊 Stockage centralisé des images médicales",
        "🌐 Standard DICOM pour l'interopérabilité",
        "📈 Volume croissant : millions d'images par an",
        "⚡ Besoin de solutions performantes et conformes",
        "🔍 Comparaison objective nécessaire"
    ])
    
    # Slide 4: Contexte médical
    print("  📄 Slide 4/25: Contexte médical")
    add_content_slide(prs, "🏥 Contexte Médical", [
        "🔹 Imagerie médicale : CT, MRI, PET, Radiographie",
        "🔹 Volumes massifs : 500+ études par jour",
        "🔹 Besoins multiples : diagnostic, recherche, enseignement",
        "🔹 Contraintes réglementaires : RGPD, HIPAA",
        "🔹 Interopérabilité entre systèmes critiques",
        "🔹 Sécurité et traçabilité essentielles"
    ])
    
    # Slide 5: Problématique
    print("  📄 Slide 5/25: Problématique")
    add_content_slide(prs, "2️⃣ Problématique", [
        "❓ Quelle solution PACS est la plus performante ?",
        "   • Orthanc : léger, open-source, rapide",
        "   • DCM4CHEE : enterprise, J2EE, scalable",
        "",
        "🔒 Comment garantir la conformité RGPD/HIPAA ?",
        "   • Anonymisation automatisée",
        "   • Traçabilité complète des modifications",
        "",
        "🎯 Comment extraire et analyser les structures RT ?",
        "   • RT-STRUCT pour planification radiothérapie",
        "   • Calcul volumes ROIs (Regions Of Interest)"
    ])
    
    # Slide 6: Objectifs
    print("  📄 Slide 6/25: Objectifs")
    add_objectives_slide(prs)
    
    # Slide 7: Indicateurs de succès
    print("  📄 Slide 7/25: Indicateurs")
    add_stats_slide(prs, "📈 Indicateurs de Succès", {
        "< 200ms": "Temps réponse",
        "18/18": "Tags HIPAA",
        "87%": "Couverture tests",
        "742": "Patients traités",
        "125 GB": "Données archivées",
        "99.7%": "Disponibilité"
    })
    
    # Slide 8: Architecture
    print("  📄 Slide 8/25: Architecture")
    add_architecture_slide(prs)
    
    # Slide 9: Stack technique
    print("  📄 Slide 9/25: Stack technique")
    add_table_slide(prs, "5️⃣ Technologies et Stack", [
        ["Couche", "Technologie", "Version", "Rôle"],
        ["Frontend", "React", "18.2", "UI/UX"],
        ["", "Vite", "5.0", "Build tool"],
        ["", "Recharts", "2.10", "Graphiques"],
        ["Backend", "FastAPI", "0.109", "API REST"],
        ["", "Python", "3.13", "Langage"],
        ["", "SQLAlchemy", "2.0", "ORM"],
        ["PACS", "Orthanc", "1.12", "PACS léger"],
        ["", "DCM4CHEE", "5.23", "Enterprise"],
        ["", "XNAT", "1.8", "Recherche"],
        ["Database", "PostgreSQL", "15", "Persistence"],
        ["Viewer", "OHIF", "3.9", "Viewer DICOM"]
    ])
    
    # Slide 10: Fonctionnalités Dashboard
    print("  📄 Slide 10/25: Dashboard")
    add_content_slide(prs, "6️⃣ Dashboard Comparatif", [
        "📊 Statistiques en temps réel",
        "   • Compteurs : Patients, Études, Séries, Instances",
        "   • Graphiques comparatifs DCM4CHEE vs Orthanc",
        "",
        "⚡ Performances",
        "   • Rafraîchissement automatique toutes les 30s",
        "   • Indicateurs de temps de réponse",
        "",
        "🎨 Visualisation",
        "   • Graphiques interactifs (Recharts)",
        "   • Filtrage par date et modalité",
        "   • Export données CSV/JSON"
    ])
    
    # Slide 11: OHIF Viewer
    print("  📄 Slide 11/25: OHIF Viewer")
    add_content_slide(prs, "7️⃣ Visualisation DICOM - OHIF", [
        "🖼️ Modes de visualisation",
        "   • 2D View : slice par slice",
        "   • MPR 3D : Axial + Sagittal + Coronal",
        "   • Volume Rendering : reconstruction 3D",
        "",
        "📏 Outils disponibles",
        "   • Mesures : distance, angle, surface, ROI",
        "   • Annotations avec export JSON",
        "   • Windowing (HU - Hounsfield Units)",
        "   • Zoom, Pan, Rotation",
        "",
        "💾 Sauvegarde automatique des annotations"
    ])
    
    # Slide 12: Anonymisation
    print("  📄 Slide 12/25: Anonymisation")
    add_table_slide(prs, "8️⃣ Anonymisation RGPD/HIPAA", [
        ["Niveau", "Tags Modifiés", "Usage", "Conformité"],
        ["Basic", "4 tags", "Usage interne", "Basique"],
        ["Partial", "12 tags + dates", "Formation, démo", "RGPD"],
        ["Full", "18 tags HIPAA", "Recherche publique", "HIPAA"]
    ])
    
    # Slide 13: Workflow Anonymisation
    print("  📄 Slide 13/25: Workflow anonymisation")
    add_content_slide(prs, "🔐 Workflow Anonymisation", [
        "1️⃣ Sélection : Choisir étude dans PACS",
        "2️⃣ Configuration : Niveau (Basic/Partial/Full)",
        "3️⃣ Prévisualisation : Voir tags avant/après",
        "4️⃣ Validation : Confirmer modifications",
        "5️⃣ Export : Upload vers XNAT",
        "",
        "🔍 Traçabilité complète",
        "   • Logs PostgreSQL avec audit trail",
        "   • Mapping original ↔ anonymisé",
        "   • Conformité 100% HIPAA (18/18 identifiants)"
    ])
    
    # Slide 14: RT-STRUCT
    print("  📄 Slide 14/25: RT-STRUCT")
    add_content_slide(prs, "9️⃣ Extraction RT-STRUCT", [
        "🎯 Analyse structures de radiothérapie",
        "   • Extraction automatique des ROIs 3D",
        "   • Calcul volumes (cm³) avec précision ±2%",
        "   • Visualisation contours et surfaces",
        "",
        "📊 Structures extraites",
        "   • GTV (Gross Tumor Volume) : tumeur visible",
        "   • PTV (Planning Target Volume) : volume cible",
        "   • OAR (Organs At Risk) : organes à protéger",
        "",
        "💾 Export JSON + PNG pour reporting",
        "⏱️ Temps extraction : 1.2s par étude"
    ])
    
    # Slide 15: Exemple RT-STRUCT
    print("  📄 Slide 15/25: Exemple RT-STRUCT")
    add_table_slide(prs, "📊 Exemple Résultats RT-STRUCT", [
        ["Structure", "Volume (cm³)", "Type", "Couleur"],
        ["GTV", "45.3", "Target", "Rouge"],
        ["PTV", "125.7", "Planning", "Orange"],
        ["Poumon Gauche", "1840", "OAR", "Bleu"],
        ["Poumon Droit", "1920", "OAR", "Bleu"],
        ["Cœur", "620", "OAR", "Vert"],
        ["Moelle Épinière", "35", "OAR", "Jaune"]
    ])
    
    # Slide 16: Comparaison Performance
    print("  📄 Slide 16/25: Comparaison performance")
    add_table_slide(prs, "10️⃣ Comparaison Performance", [
        ["Métrique", "Orthanc", "DCM4CHEE", "Gagnant"],
        ["Upload 100 DICOM", "2.3s", "4.1s", "🏆 Orthanc"],
        ["Query Patient", "45ms", "180ms", "🏆 Orthanc"],
        ["RAM Usage", "250MB", "1.8GB", "🏆 Orthanc"],
        ["Storage", "Filesystem", "PostgreSQL", "Égalité"],
        ["Scalabilité", "Moyenne", "Excellente", "🏆 DCM4CHEE"],
        ["Conformité", "Basique", "Complète", "🏆 DCM4CHEE"],
        ["Multi-site", "Limité", "Natif", "🏆 DCM4CHEE"]
    ])
    
    # Slide 17: Graphique performance
    print("  📄 Slide 17/25: Graphique performance")
    add_chart_slide(
        prs,
        "⚡ Performances Comparées",
        ["Upload", "Query", "Scalabilité"],
        [2.3, 45, 60],  # Orthanc
        [4.1, 180, 95],  # DCM4CHEE
        "Orthanc",
        "DCM4CHEE"
    )
    
    # Slide 18: Statistiques plateforme
    print("  📄 Slide 18/25: Statistiques plateforme")
    add_stats_slide(prs, "📊 Statistiques Plateforme", {
        "742": "Patients",
        "742": "Études",
        "125 GB": "Archive",
        "95ms": "Temps réponse",
        "99.7%": "Disponibilité",
        "15,000": "Req/jour"
    })
    
    # Slide 19: Tests et qualité
    print("  📄 Slide 19/25: Tests")
    add_content_slide(prs, "11️⃣ Tests et Qualité", [
        "✅ Couverture de tests : 87%",
        "✅ 29 tests passés (100% success)",
        "✅ 0 tests échoués",
        "",
        "📊 Détail par module",
        "   • models.py : 100%",
        "   • schemas.py : 100%",
        "   • database.py : 100%",
        "   • main.py : 89%",
        "   • crud.py : 90%",
        "   • sync_service.py : 63%"
    ])
    
    # Slide 20: Types de tests
    print("  📄 Slide 20/25: Types de tests")
    add_content_slide(prs, "🧪 Types de Tests Implémentés", [
        "🔹 Tests Unitaires",
        "   • Modèles de données (Patient, Study)",
        "   • Schémas Pydantic (validation)",
        "   • Opérations CRUD",
        "",
        "🔹 Tests d'Intégration",
        "   • Endpoints API (/health, /statistics)",
        "   • Workflow anonymisation",
        "   • Export XNAT",
        "",
        "🔹 Tests de Performance",
        "   • Charge 1000 requêtes/s",
        "   • Latence < 200ms"
    ])
    
    # Slide 21: Démonstration
    print("  📄 Slide 21/25: Démonstration")
    add_content_slide(prs, "12️⃣ Démonstration Live", [
        "🌐 URLs de démonstration",
        "   • Dashboard : http://localhost:8000",
        "   • OHIF Viewer : http://localhost:3001",
        "   • XNAT : http://localhost:8090",
        "",
        "🎬 Scénarios à démontrer",
        "   1. Visualisation statistiques dashboard",
        "   2. Ouverture étude DICOM dans OHIF",
        "   3. Mesures et annotations",
        "   4. Anonymisation patient",
        "   5. Export vers XNAT",
        "   6. Extraction RT-STRUCT avec volumes"
    ])
    
    # Slide 22: Perspectives court terme
    print("  📄 Slide 22/25: Perspectives")
    add_content_slide(prs, "13️⃣ Perspectives - Court Terme", [
        "📱 Application mobile (3 mois)",
        "   • Consultation à distance",
        "   • Notifications push",
        "",
        "🔔 Notifications temps réel",
        "   • WebSocket pour updates",
        "   • Alertes automatiques",
        "",
        "📊 Rapports automatisés PDF",
        "   • Génération rapports médicaux",
        "   • Templates personnalisables",
        "",
        "🌐 Support multi-langues (FR/EN/ES)"
    ])
    
    # Slide 23: Perspectives moyen terme
    print("  📄 Slide 23/25: Évolutions")
    add_content_slide(prs, "🚀 Perspectives - Moyen Terme", [
        "🤖 Intelligence Artificielle (6 mois)",
        "   • Détection automatique anomalies",
        "   • Segmentation assistée par IA",
        "",
        "📈 Analytics avancés",
        "   • Machine Learning pour prédictions",
        "   • Analyse tendances",
        "",
        "🔗 Intégration autres PACS",
        "   • Horos, RadiAnt, Synapse",
        "",
        "☁️ Déploiement cloud",
        "   • Azure, AWS, Google Cloud",
        "   • Scalabilité automatique"
    ])
    
    # Slide 24: Conclusion
    print("  📄 Slide 24/25: Conclusion")
    add_content_slide(prs, "14️⃣ Conclusion", [
        "🎓 Apprentissages clés",
        "   • Architecture microservices = Scalabilité",
        "   • Docker = Déploiement simplifié",
        "   • FastAPI + React = Stack moderne performante",
        "",
        "💡 Recommandations",
        "   • Petites structures → Orthanc (simple, rapide)",
        "   • Hôpitaux → DCM4CHEE (enterprise, scalable)",
        "   • Recherche → Orthanc + XNAT (flexible)",
        "",
        "🏆 Résultats obtenus",
        "   ✅ Plateforme fonctionnelle et performante",
        "   ✅ Conformité RGPD/HIPAA garantie",
        "   ✅ Documentation complète (70 pages)"
    ])
    
    # Slide 25: Merci
    print("  📄 Slide 25/25: Merci")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(102, 126, 234)
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Merci ! 🙏"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(6), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions ?"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(36)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_file = "C:\\Users\\awati\\Desktop\\pacs\\Presentation_PACS_Professionnelle.pptx"
    prs.save(output_file)
    
    print(f"\n✅ Présentation générée avec succès !")
    print(f"📁 Fichier : {output_file}")
    print(f"📊 Nombre de slides : {len(prs.slides)}")
    print(f"\n🚀 Pour ouvrir : Start-Process '{output_file}'")
    
    return output_file

if __name__ == "__main__":
    generate_presentation()
